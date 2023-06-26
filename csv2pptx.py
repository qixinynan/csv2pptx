from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
import csv


def add_text_slide(_presentation, text, font_size: int, width_num: int, height_num: int, top_num: int, left_num: int, font_name: str):
    title_slide_layout = _presentation.slide_layouts[6]
    slide = _presentation.slides.add_slide(title_slide_layout)

    width = Pt(width_num)
    height = Pt(height_num)
    top = Pt(top_num)
    left = Pt(left_num)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    # 设置文本框中的文本
    text_frame.text = text

    # 设置文本框中文本的样式
    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(font_size)
    paragraph.font.name = font_name
    paragraph.font.bold = True
    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER


def generate_presentation_by_array(array: [str], font_size: int, width_num: int, height_num: int,
                                   top_num: int, left_num: int, font_name: str):
    _presentation = Presentation()

    for item in array:
        add_text_slide(_presentation, item, font_size, width_num, height_num, top_num, left_num, font_name)
    return _presentation


def read_csv(file, start_row: int, start_column: int) -> [str]:
    csv_reader = csv.reader(open(file))
    _data: [str] = []
    for i in range(start_row - 1):
        next(csv_reader)
    for row in csv_reader:
        item = row[start_column - 1]
        _data.append(item)
    return _data


# data = read_csv("input.csv", 3, 2)
# sample = random.sample(data, 50)
# presentation = generate_presentation_by_array(sample)
# presentation.save("output.pptx")
# print(sample)
